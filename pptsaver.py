import os
from pptx import Presentation
from google.adk.tools import FunctionTool
from utils.config import OUTPUT_DIRECTORY


def create_pptx_file(title: str, slides_content: list[dict], filename: str ="output_ppt"):
    '''
    Creates a Powerpoint file
    '''
    try:
        os.makedirs(OUTPUT_DIRECTORY,exist_ok=True)
        filepath=os.path.join(OUTPUT_DIRECTORY,filename)

        ppt=Presentation()
        # adding title slide
        title_slide_layout=ppt.slide_layouts[0]
        slide =ppt.slides.add_slide(title_slide_layout)
        slide.shapes.title.text=title
        
        #adding the content slides

        main_slide_layout=ppt.slide_layouts[1]
        count =1
        for slide_data in slides_content:
            slide= ppt.slides.add_slide(main_slide_layout)
            slide.shapes.title.text=slide_data.get('title','Slide No '+str(count))

            count+=1

            body_shape=slide.shapes.placeholders[1]
            text_frame=body_shape.text_frame
            text_frame.text=slide_data.get('content','Slide has no content')
        
        # Save ppt to Output Directory
        ppt.save(filepath)

        return f" PPT saved to destination {filepath}"
    #exception behaviour to be tackled by exception clause
    except Exception as e:
        return f"PPT cannot be saved due to error: {str(e)}"
    
PPTSaverTool=FunctionTool(create_pptx_file)
    
     




    
